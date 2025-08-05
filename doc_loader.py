from paddleocr import PaddleOCR
from PIL import Image
import fitz  # PyMuPDF
import io
import os
import numpy as np
import pandas as pd
import openpyxl
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation

# Initialize PaddleOCR with updated parameter
ocr = PaddleOCR(use_textline_orientation=True, lang='en')

def load_and_chunk_pdf(file_path, chunk_size=1000, overlap=100):
    print(f"🔍 Processing file: {file_path}")
    print(f"📁 File exists: {os.path.exists(file_path)}")
    
    if not os.path.exists(file_path):
        print("❌ File not found!")
        return
    
    chunks, metadata = [], []
    
    try:
        doc = fitz.open(file_path)
        
        filename = os.path.basename(file_path)
        
        for page_num in range(len(doc)):
    
            page = doc[page_num]
            text = page.get_text()

            image_list = page.get_images()
            image_texts = []
            
            for img_index, img in enumerate(image_list):
                try:
                    print(f"   Processing image {img_index + 1}...")
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    # Skip if not RGB or GRAY
                    if pix.n - pix.alpha < 4:
                        img_data = pix.tobytes("png")
                        pil_image = Image.open(io.BytesIO(img_data))
                        img_array = np.array(pil_image)
                        ocr_result = ocr.ocr(img_array)
                        extracted_text = []
                        if ocr_result and len(ocr_result) > 0:
                            result_data = ocr_result[0]
                            
                            if isinstance(result_data, dict):
                                rec_texts = result_data.get('rec_texts', [])
                                rec_scores = result_data.get('rec_scores', [])
                                                
                                for text_content, confidence in zip(rec_texts, rec_scores):
                                    if confidence > 0.3 and text_content.strip():
                                        extracted_text.append(text_content.strip())
                                        print(f"   ✅ OCR Text: '{text_content}' (confidence: {confidence:.3f})")
                            
                            elif isinstance(result_data, list):
                                # Handle standard list format (fallback)
                                for line in result_data:
                                    if line and len(line) >= 2:
                                        text_content = line[1][0]
                                        confidence = line[1][1]
                                        
                                        if confidence > 0.3 and text_content.strip():
                                            extracted_text.append(text_content.strip())
                                            print(f"   ✅ OCR Text: '{text_content}' (confidence: {confidence:.3f})")
                        
                        if extracted_text:
                            combined_text = " ".join(extracted_text)
                            image_texts.append(f"[IMAGE_TEXT]: {combined_text}")
                        else:
                            print(f"   ❌ No text extracted from image {img_index + 1}")
                    
                    pix = None
                    
                except Exception as e:
                    print(f"   ❌ OCR Error on image {img_index + 1}: {e}")
                    continue
            
            # Combine all text
            full_text = text
            if image_texts:
                full_text += "\n\n" + "\n".join(image_texts)
            
            print(f"📋 Combined text length: {len(full_text.strip())}")
            
            # Create chunks
            page_chunks = 0
            if full_text.strip():
                for i in range(0, len(full_text), chunk_size - overlap):
                    chunk = full_text[i:i+chunk_size]
                    if chunk.strip():
                        chunks.append(chunk)
                        metadata.append({
                            "page": page_num + 1,
                            "source": filename,
                            "text": chunk,
                            "has_images": len(image_texts) > 0,
                            "image_count": len(image_texts)
                        })
                        page_chunks += 1
            
        
        doc.close()
        
        return chunks, metadata
        
    except Exception as e:
        print(f"❌ Error processing PDF: {e}")
        return [], []


def load_and_chunk_excel(file_path, rows_per_chunk=5):
    filename = os.path.basename(file_path)
    chunks, metadata = [], []

    with pd.ExcelFile(file_path) as xls:
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)

            total_rows = len(df)
            if total_rows == 0:
                continue

            # Iterate over the sheet in row-size steps
            for start in range(0, total_rows, rows_per_chunk):
                end = min(start + rows_per_chunk, total_rows)

                chunk_df = df.iloc[start:end]

                # Convert to string, replace NaNs, and concatenate
                chunk_text = (
                    chunk_df.fillna("")            # replace NaNs
                            .astype(str)
                            .agg(" ".join, axis=1)  # join columns per row
                            .str.cat(sep=" ")       # join rows into one string
                            .strip()
                )

                # Skip completely empty chunks
                if not chunk_text:
                    continue

                chunks.append(chunk_text)
                metadata.append({
                    "source": filename,
                    "sheet": sheet_name,
                    "start_row": start + 1,  # +1 for human-readable numbering
                    "end_row": end,
                    "text": chunk_text
                })

    return chunks, metadata



def load_and_chunk_ppt(file_path: str, chars_per_chunk: int = 1000, overlap: int = 100):
    print(f"🔍 Processing PowerPoint: {file_path}")
    
    if not os.path.exists(file_path):
        print("❌ File not found!")
        return [], []

    filename = os.path.basename(file_path)
    chunks, metadata = [], []

    try:
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            print(f"   Processing slide {slide_idx}...")
            
            # ── 1️⃣ Extract text from shapes
            text_parts = []
            shape_count = 0
            for shape in slide.shapes:
                shape_count += 1
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    text_parts.append(text)
                    print(f"     📝 Shape {shape_count}: {text[:50]}...")
            
            slide_text = "\n".join(text_parts)
            print(f"     📋 Slide {slide_idx} text length: {len(slide_text)} characters")
            print(f"     📋 Total shapes on slide: {shape_count}")
            
            # ── 2️⃣ OCR every picture on the slide
            image_texts = []
            picture_count = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1
                    try:
                        print(f"     🖼️ Processing image {picture_count}...")
                        
                        # Extract image data
                        img_bytes = shape.image.blob
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        print(f"        Image size: {pil_img.size}, mode: {pil_img.mode}")
                        
                        # Convert to numpy array for OCR
                        img_array = np.array(pil_img)
                        print(f"        Array shape: {img_array.shape}")
                        
                        # Run OCR
                        ocr_result = ocr.ocr(img_array)
                        extracted = []
                        
                        if ocr_result and len(ocr_result) > 0:
                            result_data = ocr_result[0]
                            
                            # Handle new dict format (Paddle >= 2.6)
                            if isinstance(result_data, dict):
                                rec_texts = result_data.get('rec_texts', [])
                                rec_scores = result_data.get('rec_scores', [])
                                
                                for txt, conf in zip(rec_texts, rec_scores):
                                    if conf > 0.3 and txt.strip():
                                        extracted.append(txt.strip())
                                        print(f"        ✅ OCR Text: '{txt}' (confidence: {conf:.3f})")
                            
                            # Handle list format (older versions)
                            elif isinstance(result_data, list):
                                for line in result_data:
                                    if line and len(line) >= 2:
                                        txt, conf = line[1]
                                        if conf > 0.3 and txt.strip():
                                            extracted.append(txt.strip())
                                            print(f"        ✅ OCR Text: '{txt}' (confidence: {conf:.3f})")
                        
                        if extracted:
                            combined_text = "[IMAGE_TEXT]: " + " ".join(extracted)
                            image_texts.append(combined_text)
                            print(f"        📝 Combined OCR: {combined_text[:100]}...")
                        else:
                            print(f"        ❌ No text extracted from image {picture_count}")
                            
                    except Exception as e:
                        print(f"        ❌ OCR error on image {picture_count}: {e}")
                        continue
            
            print(f"     🖼️ Found {picture_count} images on slide {slide_idx}")
            
            # ── 3️⃣ Combine text sources
            full_text = slide_text
            if image_texts:
                full_text += "\n\n" + "\n".join(image_texts)
            
            print(f"     📄 Combined text length: {len(full_text)} characters")
            
            if not full_text.strip():
                print(f"     ⚠️ Slide {slide_idx} is completely empty, skipping...")
                continue
            
            # ── 4️⃣ Create chunks with overlap
            pos = 0
            chunk_count = 0
            while pos < len(full_text):
                chunk = full_text[pos:pos + chars_per_chunk].strip()
                if chunk:
                    chunks.append(chunk)
                    metadata.append({
                        "source": filename,
                        "slide": slide_idx,
                        "text": chunk,
                        "has_images": bool(image_texts),
                        "image_count": len(image_texts)
                    })
                    chunk_count += 1
                    print(f"        📦 Chunk {chunk_count}: {len(chunk)} chars")
                pos += chars_per_chunk - overlap
            
            print(f"     ✅ Created {chunk_count} chunks from slide {slide_idx}")
        
        print(f"🎉 Total chunks created: {len(chunks)}")
        print(f"🎉 Total metadata entries: {len(metadata)}")
        
        return chunks, metadata
        
    except Exception as e:
        print(f"❌ Error processing PPTX: {e}")
        import traceback
        traceback.print_exc()
        return [], []
